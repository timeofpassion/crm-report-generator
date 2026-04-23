import io
from typing import Dict, Any, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.oxml.ns import qn
import pandas as pd
import numpy as np

from modules import chart_utils as cu


# ── 색상 ──────────────────────────────────────────────────────────────────────
CP  = RGBColor(0x1A, 0x3A, 0x6B)  # Primary  (navy)
CS  = RGBColor(0x2D, 0x7D, 0xD2)  # Secondary (blue)
CA  = RGBColor(0xE8, 0xA0, 0x20)  # Accent   (amber)
CL  = RGBColor(0xEB, 0xF3, 0xFB)  # Light bg
CW  = RGBColor(0xFF, 0xFF, 0xFF)  # White
CD  = RGBColor(0x1C, 0x2B, 0x3A)  # Dark text
CG  = RGBColor(0x6B, 0x72, 0x80)  # Gray
CG2 = RGBColor(0xF3, 0xF4, 0xF6)  # Light gray bg

SW = Inches(13.333)
SH = Inches(7.5)

FONT = '맑은 고딕'


def _fmt(v):
    v = float(v)
    if v >= 100_000_000:
        return f'{v/100_000_000:.1f}억원'
    elif v >= 10_000:
        return f'{v/10_000:.0f}만원'
    return f'{v:,.0f}원'


def _pct(v):
    return f'{float(v):.1f}%'


class PPTGenerator:
    def __init__(self, hospital_config: dict):
        self.cfg = hospital_config

    # ── PUBLIC ─────────────────────────────────────────────────────────────────

    def generate(self, analysis: Dict[str, Any], report_cfg: dict) -> bytes:
        prs = Presentation()
        prs.slide_width  = SW
        prs.slide_height = SH
        blank = prs.slide_layouts[6]

        self._slide_cover(prs, blank, report_cfg)
        self._slide_executive_summary(prs, blank, analysis, report_cfg)
        self._slide_category_sales(prs, blank, analysis, report_cfg)
        self._slide_reclassify(prs, blank, analysis, report_cfg)
        self._slide_cross_selling(prs, blank, analysis, report_cfg)
        self._slide_doctor_sales(prs, blank, analysis, report_cfg)
        self._slide_channel_sales(prs, blank, analysis, report_cfg)
        self._slide_channel_x_category(prs, blank, analysis, report_cfg)
        self._slide_nationality(prs, blank, analysis, report_cfg)
        self._slide_patient_spending(prs, blank, analysis, report_cfg)
        self._slide_trends(prs, blank, analysis, report_cfg)
        self._slide_insights(prs, blank, analysis, report_cfg)
        self._slide_appendix_cover(prs, blank, report_cfg)
        self._slide_appendix_category_detail(prs, blank, analysis, report_cfg)
        self._slide_appendix_lifting(prs, blank, analysis, report_cfg)
        self._slide_appendix_vip(prs, blank, analysis, report_cfg)
        self._slide_appendix_misc(prs, blank, analysis, report_cfg)

        out = io.BytesIO()
        prs.save(out)
        out.seek(0)
        return out.getvalue()

    # ── HELPERS ────────────────────────────────────────────────────────────────

    def _add_slide(self, prs, blank):
        return prs.slides.add_slide(blank)

    def _rect(self, slide, left, top, w, h, fill: RGBColor, line: Optional[RGBColor] = None):
        shape = slide.shapes.add_shape(1, left, top, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        if line:
            shape.line.color.rgb = line
        else:
            shape.line.fill.background()
        return shape

    def _text(self, slide, text, left, top, w, h,
              size=11, bold=False, color: RGBColor = CD,
              align=PP_ALIGN.LEFT, italic=False, wrap=True):
        box = slide.shapes.add_textbox(left, top, w, h)
        tf  = box.text_frame
        tf.word_wrap = wrap
        p   = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.italic = italic
        run.font.name  = FONT
        run.font.color.rgb = color
        return box

    def _header(self, slide, title: str, page: int, report_cfg: dict):
        # Header bar
        self._rect(slide, 0, 0, SW, Inches(0.85), CP)
        # Title
        self._text(slide, title,
                   Inches(0.45), Inches(0.12), Inches(10), Inches(0.65),
                   size=20, bold=True, color=CW)
        # Page number
        self._text(slide, f'{page} / 17',
                   SW - Inches(1.1), Inches(0.22), Inches(0.9), Inches(0.4),
                   size=10, color=CW, align=PP_ALIGN.RIGHT)
        # Footer
        footer = f"{report_cfg['hospital_name']} | {report_cfg['year']}년 {report_cfg['month']}월 | {report_cfg['team_name']}"
        self._text(slide, footer,
                   Inches(0.4), SH - Inches(0.35), SW - Inches(0.8), Inches(0.3),
                   size=7, color=CG)

    def _kpi_box(self, slide, left, top, w, h, label, value, sub=''):
        self._rect(slide, left, top, w, h, CL)
        self._text(slide, label, left + Inches(0.15), top + Inches(0.1), w - Inches(0.3), Inches(0.3),
                   size=9, color=CG)
        self._text(slide, value, left + Inches(0.1), top + Inches(0.38), w - Inches(0.2), Inches(0.55),
                   size=18, bold=True, color=CP)
        if sub:
            self._text(slide, sub, left + Inches(0.15), top + Inches(0.88), w - Inches(0.3), Inches(0.25),
                       size=8, color=CG, italic=True)

    def _img(self, slide, buf, left, top, w, h):
        if buf:
            buf.seek(0)
            slide.shapes.add_picture(buf, left, top, w, h)

    def _table_text(self, slide, headers, rows, left, top, w, h,
                    header_color: RGBColor = CP, row_colors=None):
        n_cols = len(headers)
        n_rows = len(rows) + 1
        tbl = slide.shapes.add_table(n_rows, n_cols, left, top, w, h).table

        col_w = w // n_cols
        for i in range(n_cols):
            tbl.columns[i].width = col_w

        for ci, hdr in enumerate(headers):
            cell = tbl.cell(0, ci)
            cell.text = str(hdr)
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_color
            p  = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            run.font.bold  = True
            run.font.size  = Pt(8)
            run.font.color.rgb = CW
            run.font.name  = FONT

        for ri, row in enumerate(rows):
            bg = row_colors[ri % len(row_colors)] if row_colors else (CW if ri % 2 == 0 else CG2)
            for ci, val in enumerate(row):
                cell = tbl.cell(ri + 1, ci)
                cell.text = str(val)
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg
                p  = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.runs[0] if p.runs else p.add_run()
                run.font.size  = Pt(8)
                run.font.color.rgb = CD
                run.font.name  = FONT

    # ── SLIDES ─────────────────────────────────────────────────────────────────

    def _slide_cover(self, prs, blank, rc):
        slide = self._add_slide(prs, blank)
        # Full background
        self._rect(slide, 0, 0, SW, SH, CP)
        # Accent bar
        self._rect(slide, 0, SH - Inches(1.5), SW, Inches(4), CS)
        self._rect(slide, 0, SH - Inches(1.2), SW, Inches(0.12), CA)

        # Titles
        self._text(slide, rc['hospital_name'],
                   Inches(1), Inches(1.8), SW - Inches(2), Inches(1.2),
                   size=40, bold=True, color=CW, align=PP_ALIGN.CENTER)
        self._text(slide, '매출분석 보고서',
                   Inches(1), Inches(3.0), SW - Inches(2), Inches(0.9),
                   size=28, color=CW, align=PP_ALIGN.CENTER)
        self._text(slide, f"{rc['year']}년 {rc['month']}월",
                   Inches(1), Inches(3.9), SW - Inches(2), Inches(0.7),
                   size=20, color=CA, align=PP_ALIGN.CENTER, bold=True)
        self._text(slide, rc['team_name'],
                   Inches(1), SH - Inches(1.0), SW - Inches(2), Inches(0.5),
                   size=12, color=CG2, align=PP_ALIGN.CENTER)

    def _slide_executive_summary(self, prs, blank, analysis, rc):
        slide = self._add_slide(prs, blank)
        self._header(slide, 'Executive Summary', 2, rc)

        kpis = analysis.get('kpis', {})
        # 4 KPI boxes
        bw, bh = Inches(2.9), Inches(1.25)
        gap = Inches(0.15)
        start_x = Inches(0.4)
        y = Inches(1.0)

        boxes = [
            ('총 매출', _fmt(kpis.get('total_revenue', 0)), ''),
            ('총 방문건수', f"{kpis.get('total_visits', 0):,}건", ''),
            ('총 환자수', f"{kpis.get('unique_patients', 0):,}명", ''),
            ('평균 객단가', _fmt(kpis.get('avg_per_patient', 0)), '(환자 1인당)'),
        ]
        for i, (lbl, val, sub) in enumerate(boxes):
            self._kpi_box(slide, start_x + i * (bw + gap), y, bw, bh, lbl, val, sub)

        # Key findings
        self._rect(slide, Inches(0.4), Inches(2.5), SW - Inches(0.8), Inches(0.35), CP)
        self._text(slide, '핵심 발견 (Key Findings)',
                   Inches(0.55), Inches(2.55), Inches(6), Inches(0.28),
                   size=11, bold=True, color=CW)

        findings = analysis.get('key_findings', [])
        for i, f in enumerate(findings[:4]):
            fy = Inches(3.0) + i * Inches(0.72)
            self._rect(slide, Inches(0.4), fy, Inches(0.35), Inches(0.35), CS)
            self._text(slide, str(i + 1), Inches(0.44), fy + Inches(0.04),
                       Inches(0.3), Inches(0.3), size=12, bold=True, color=CW, align=PP_ALIGN.CENTER)
            self._text(slide, f, Inches(0.85), fy + Inches(0.04),
                       SW - Inches(1.3), Inches(0.6), size=11, color=CD)

    def _slide_category_sales(self, prs, blank, analysis, rc):
        cat = analysis.get('category_sales')
        slide = self._add_slide(prs, blank)
        self._header(slide, '시술 카테고리별 매출 분석', 3, rc)

        if cat is None or cat.empty:
            self._text(slide, '데이터 없음', Inches(5), Inches(3.5), Inches(3), Inches(1),
                       size=14, color=CG, align=PP_ALIGN.CENTER)
            return

        labels = cat['category'].tolist()
        values = cat['revenue'].tolist()

        buf = cu.donut_chart(labels, values, figsize=(6, 5))
        self._img(slide, buf, Inches(0.3), Inches(0.9), Inches(6.5), Inches(5.5))

        # Table
        headers = ['카테고리', '매출', '건수', '비율']
        rows = [[r['category'], _fmt(r['revenue']), f"{int(r['count']):,}", f"{r['pct']:.1f}%"]
                for _, r in cat.iterrows()]
        self._table_text(slide, headers, rows[:8],
                         Inches(7.0), Inches(1.0), Inches(5.9), Inches(5.5))

    def _slide_reclassify(self, prs, blank, analysis, rc):
        slide = self._add_slide(prs, blank)
        self._header(slide, '기존 분류 vs 재분류 비교', 4, rc)

        orig = analysis.get('original_categories')
        reclassify = analysis.get('reclassify_compare')
        cat = analysis.get('category_sales')

        if reclassify is not None and not reclassify.empty:
            df = reclassify
            headers = ['구분', 'CRM 원래 분류 매출', '재분류 후 매출', '차이']
            rows = []
            for _, r in df.iterrows():
                diff = r.get('new_revenue', 0) - r.get('orig_revenue', 0)
                rows.append([
                    r.get('original_category', '-'),
                    _fmt(r.get('orig_revenue', 0)),
                    _fmt(r.get('new_revenue', 0)),
                    ('+' if diff >= 0 else '') + _fmt(abs(diff)),
                ])
            self._table_text(slide, headers, rows[:12],
                             Inches(0.4), Inches(1.0), SW - Inches(0.8), Inches(5.5))
        elif cat is not None and not cat.empty:
            labels = cat['category'].tolist()
            values = cat['revenue'].tolist()
            buf = cu.hbar_chart(labels, values, title='재분류 기준 카테고리별 매출', figsize=(9, 5))
            self._img(slide, buf, Inches(1.5), Inches(1.0), Inches(10), Inches(5.5))
        else:
            self._text(slide, '오더판매내역 파일을 업로드하면 기존/재분류 비교가 가능합니다.',
                       Inches(2), Inches(3.5), Inches(9), Inches(1),
                       size=13, color=CG, align=PP_ALIGN.CENTER)

    def _slide_cross_selling(self, prs, blank, analysis, rc):
        slide = self._add_slide(prs, blank)
        self._header(slide, '복합시술 (Cross-selling) 분석', 5, rc)

        cs = analysis.get('cross_selling')
        if not cs:
            self._text(slide, '오더판매내역 데이터가 필요합니다.',
                       Inches(3), Inches(3.5), Inches(7), Inches(1),
                       size=13, color=CG, align=PP_ALIGN.CENTER)
            return

        # KPI boxes
        self._kpi_box(slide, Inches(0.5), Inches(1.0), Inches(3.5), Inches(1.3),
                      '복합시술 비율', f"{cs['cross_sell_rate']:.1f}%", '2개 이상 시술 환자')
        self._kpi_box(slide, Inches(4.3), Inches(1.0), Inches(3.5), Inches(1.3),
                      '객단가 배수', f"{cs['avg_multiplier']:.2f}x", '단일 대비')

        # Top combinations table
        combos = cs.get('top_combinations')
        if combos is not None and not combos.empty:
            self._rect(slide, Inches(0.5), Inches(2.6), Inches(7.5), Inches(0.38), CP)
            self._text(slide, 'TOP 복합시술 조합', Inches(0.65), Inches(2.65),
                       Inches(6), Inches(0.3), size=11, bold=True, color=CW)

            headers = ['시술 조합', '횟수']
            rows = [[r['조합'], f"{int(r['횟수']):,}건"] for _, r in combos.iterrows()]
            self._table_text(slide, headers, rows,
                             Inches(0.5), Inches(3.1), Inches(7.5), Inches(3.2))

        # Insight
        self._rect(slide, Inches(8.3), Inches(1.0), Inches(4.7), Inches(5.5), CL)
        self._text(slide, '💡 인사이트',
                   Inches(8.5), Inches(1.15), Inches(4.3), Inches(0.4),
                   size=11, bold=True, color=CP)
        insights = [
            f"복합시술 환자가 전체 {cs['cross_sell_rate']:.1f}%를 차지합니다.",
            f"복합시술 고객은 단일 고객 대비 {cs['avg_multiplier']:.1f}배 높은 객단가를 보입니다.",
            "TOP 조합을 패키지 상품화하면 추가 매출 창출이 가능합니다.",
        ]
        for i, ins in enumerate(insights):
            self._text(slide, f'• {ins}',
                       Inches(8.5), Inches(1.65) + i * Inches(0.75), Inches(4.3), Inches(0.7),
                       size=10, color=CD, wrap=True)

    def _slide_doctor_sales(self, prs, blank, analysis, rc):
        slide = self._add_slide(prs, blank)
        self._header(slide, '진료의별 매출 비교', 6, rc)

        doc = analysis.get('doctor_sales')
        if doc is None or doc.empty:
            self._text(slide, '진료의 데이터 없음', Inches(5), Inches(3.5), Inches(3), Inches(1),
                       size=14, color=CG, align=PP_ALIGN.CENTER)
            return

        buf = cu.grouped_bar_chart(doc, 'doctor', 'category', 'revenue',
                                   title='진료의별 카테고리 매출', figsize=(9, 5))
        self._img(slide, buf, Inches(0.3), Inches(0.9), Inches(9), Inches(5.5))

        # Summary table
        doc_total = doc.groupby('doctor')['revenue'].sum().sort_values(ascending=False).reset_index()
        headers = ['진료의', '총 매출']
        rows = [[r['doctor'], _fmt(r['revenue'])] for _, r in doc_total.iterrows()]
        self._table_text(slide, headers, rows[:8],
                         Inches(9.5), Inches(1.0), Inches(3.5), Inches(5.5))

    def _slide_channel_sales(self, prs, blank, analysis, rc):
        slide = self._add_slide(prs, blank)
        self._header(slide, '유입경로별 매출 분석', 7, rc)

        ch = analysis.get('channel_sales')
        if ch is None or ch.empty:
            self._text(slide, '유입경로 데이터 없음', Inches(5), Inches(3.5), Inches(3), Inches(1),
                       size=14, color=CG, align=PP_ALIGN.CENTER)
            return

        labels = ch['channel'].tolist()[:10]
        values = ch['revenue'].tolist()[:10]

        buf = cu.hbar_chart(labels, values, title='채널별 매출', figsize=(7, 5))
        self._img(slide, buf, Inches(0.3), Inches(0.9), Inches(7.5), Inches(5.5))

        # Top 5 avg table
        ch_sorted = ch.sort_values('avg', ascending=False).head(5)
        headers = ['채널', '평균 객단가', '방문수']
        rows = [[r['channel'], _fmt(r.get('avg', 0)), f"{int(r.get('count', 0)):,}건"]
                for _, r in ch_sorted.iterrows()]
        self._rect(slide, Inches(8.0), Inches(0.9), Inches(5.0), Inches(0.38), CS)
        self._text(slide, '채널별 평균 객단가 TOP 5',
                   Inches(8.15), Inches(0.97), Inches(4.5), Inches(0.3),
                   size=10, bold=True, color=CW)
        self._table_text(slide, headers, rows,
                         Inches(8.0), Inches(1.35), Inches(5.0), Inches(2.5))

    def _slide_channel_x_category(self, prs, blank, analysis, rc):
        slide = self._add_slide(prs, blank)
        self._header(slide, '유입경로 × 시술 카테고리 교차분석', 8, rc)

        pivot = analysis.get('channel_x_category')
        if pivot is None or pivot.empty:
            self._text(slide, '데이터 없음', Inches(5), Inches(3.5), Inches(3), Inches(1),
                       size=14, color=CG, align=PP_ALIGN.CENTER)
            return

        # Show top 8 channels
        top_channels = analysis.get('channel_sales', pd.DataFrame())
        if not top_channels.empty and 'channel' in top_channels.columns:
            top_ch_names = top_channels.head(8)['channel'].tolist()
            pivot = pivot.reindex([c for c in top_ch_names if c in pivot.index])

        buf = cu.heatmap(pivot.fillna(0), title='채널별 카테고리 매출 (단위: 원)', figsize=(9, 5))
        self._img(slide, buf, Inches(0.5), Inches(0.9), Inches(12.3), Inches(5.5))

    def _slide_nationality(self, prs, blank, analysis, rc):
        slide = self._add_slide(prs, blank)
        self._header(slide, '국적별 매출 & 선호 시술 분석', 9, rc)

        nat = analysis.get('nationality_sales')
        if nat is None or nat.empty:
            self._text(slide, '국적 데이터 없음', Inches(5), Inches(3.5), Inches(3), Inches(1),
                       size=14, color=CG, align=PP_ALIGN.CENTER)
            return

        # Domestic vs foreign pie
        nat_display = nat.copy()
        domestic = nat_display[nat_display['nationality'].isin(['한국', '내국인', 'KOR', 'Korean', '한국인'])]
        foreign  = nat_display[~nat_display['nationality'].isin(['한국', '내국인', 'KOR', 'Korean', '한국인'])]

        dom_rev = domestic['revenue'].sum()
        for_rev = foreign['revenue'].sum()

        if dom_rev + for_rev > 0:
            buf = cu.donut_chart(['내국인', '외국인'], [dom_rev, for_rev], figsize=(5, 4))
            self._img(slide, buf, Inches(0.3), Inches(0.9), Inches(5.5), Inches(5.5))

        # Nationality table
        headers = ['국적', '매출', '방문수', '평균 객단가']
        rows = [[r['nationality'], _fmt(r['revenue']),
                 f"{int(r.get('count', 0)):,}건", _fmt(r.get('avg', 0))]
                for _, r in nat.head(10).iterrows()]
        self._table_text(slide, headers, rows,
                         Inches(6.0), Inches(1.0), Inches(6.9), Inches(5.5))

    def _slide_patient_spending(self, prs, blank, analysis, rc):
        slide = self._add_slide(prs, blank)
        self._header(slide, '환자 객단가 분포 & VIP 분석', 10, rc)

        pt = analysis.get('patient_spending')
        if pt is None or pt.empty:
            self._text(slide, '데이터 없음', Inches(5), Inches(3.5), Inches(3), Inches(1),
                       size=14, color=CG, align=PP_ALIGN.CENTER)
            return

        totals = pt['total'].values
        buf = cu.histogram(totals, title='환자 객단가 분포', figsize=(7, 4))
        self._img(slide, buf, Inches(0.3), Inches(0.9), Inches(7.5), Inches(5.5))

        # Stats
        stats = [
            ('평균 객단가',   _fmt(np.mean(totals))),
            ('중앙값 (Median)', _fmt(np.median(totals))),
            ('상위 10% (P90)',  _fmt(np.percentile(totals, 90))),
            ('최고 객단가',    _fmt(np.max(totals))),
            ('최저 객단가',    _fmt(np.min(totals))),
        ]
        self._rect(slide, Inches(8.0), Inches(0.9), Inches(5.0), Inches(0.38), CP)
        self._text(slide, '핵심 통계', Inches(8.15), Inches(0.97),
                   Inches(4.5), Inches(0.3), size=10, bold=True, color=CW)

        for i, (lbl, val) in enumerate(stats):
            fy = Inches(1.4) + i * Inches(0.7)
            self._rect(slide, Inches(8.0), fy, Inches(5.0), Inches(0.6), CG2)
            self._text(slide, lbl, Inches(8.15), fy + Inches(0.05),
                       Inches(3.0), Inches(0.28), size=8, color=CG)
            self._text(slide, val, Inches(8.15), fy + Inches(0.28),
                       Inches(4.7), Inches(0.28), size=13, bold=True, color=CP)

    def _slide_trends(self, prs, blank, analysis, rc):
        slide = self._add_slide(prs, blank)
        self._header(slide, '일별/요일별 매출 추이', 11, rc)

        daily = analysis.get('daily_trends')
        wd    = analysis.get('weekday_trends')

        if daily is None or daily.empty:
            self._text(slide, '날짜 데이터 없음', Inches(5), Inches(3.5), Inches(3), Inches(1),
                       size=14, color=CG, align=PP_ALIGN.CENTER)
            return

        buf_line = cu.line_chart(
            list(range(len(daily))), daily['revenue'].tolist(),
            title=f"{rc['year']}년 {rc['month']}월 일별 매출", figsize=(9, 3.5)
        )
        self._img(slide, buf_line, Inches(0.3), Inches(0.9), Inches(12.7), Inches(3.5))

        if wd is not None and not wd.empty:
            buf_wd = cu.weekday_bar_chart(
                wd['weekday_name'].tolist(), wd['revenue'].tolist(),
                title='요일별 매출', figsize=(6, 3)
            )
            self._img(slide, buf_wd, Inches(0.3), Inches(4.4), Inches(12.7), Inches(2.8))

    def _slide_insights(self, prs, blank, analysis, rc):
        slide = self._add_slide(prs, blank)
        self._header(slide, '종합 인사이트 & 전략 제언', 12, rc)

        insights = analysis.get('insights', {})
        perspectives = [
            ('원장님 관점', insights.get('원장님 관점', []), CS),
            ('마케팅 관점', insights.get('마케팅 관점', []), CA),
            ('운영 관점',   insights.get('운영 관점', []),   RGBColor(0x27, 0xAE, 0x60)),
        ]

        col_w = Inches(4.1)
        gap   = Inches(0.25)
        x     = Inches(0.3)

        for i, (title, items, color) in enumerate(perspectives):
            cx = x + i * (col_w + gap)
            self._rect(slide, cx, Inches(0.9), col_w, Inches(0.45), color)
            self._text(slide, title, cx + Inches(0.1), Inches(0.96),
                       col_w - Inches(0.2), Inches(0.35),
                       size=11, bold=True, color=CW, align=PP_ALIGN.CENTER)
            self._rect(slide, cx, Inches(1.35), col_w, Inches(5.6), CG2)

            for j, item in enumerate(items[:4]):
                iy = Inches(1.5) + j * Inches(1.25)
                self._rect(slide, cx + Inches(0.1), iy, Inches(0.25), Inches(0.25), color)
                self._text(slide, item,
                           cx + Inches(0.45), iy - Inches(0.02),
                           col_w - Inches(0.6), Inches(1.1),
                           size=9, color=CD, wrap=True)

    def _slide_appendix_cover(self, prs, blank, rc):
        slide = self._add_slide(prs, blank)
        self._rect(slide, 0, 0, SW, SH, CG2)
        self._rect(slide, 0, Inches(2.5), SW, Inches(2.5), CS)
        self._text(slide, '부 록', Inches(0), Inches(2.85), SW, Inches(1.2),
                   size=36, bold=True, color=CW, align=PP_ALIGN.CENTER)
        self._text(slide, 'Appendix — 상세 데이터',
                   Inches(0), Inches(4.0), SW, Inches(0.6),
                   size=16, color=CW, align=PP_ALIGN.CENTER)
        items = ['카테고리 세부 내역', '리프팅 기기별 분석', 'VIP 환자 목록', '기타(미분류) 상세']
        for i, it in enumerate(items):
            self._text(slide, f'• {it}',
                       Inches(4.5), Inches(4.9) + i * Inches(0.42), Inches(5), Inches(0.4),
                       size=11, color=CP)

    def _slide_appendix_category_detail(self, prs, blank, analysis, rc):
        slide = self._add_slide(prs, blank)
        self._header(slide, '[부록] 카테고리별 상세 매출', 14, rc)

        cat = analysis.get('category_sales')
        if cat is None or cat.empty:
            self._text(slide, '데이터 없음', Inches(5), Inches(3.5), Inches(3), Inches(1),
                       size=14, color=CG, align=PP_ALIGN.CENTER)
            return

        labels = cat['category'].tolist()
        values = cat['revenue'].tolist()
        buf = cu.bar_chart(labels, values, title='카테고리별 매출', figsize=(9, 4.5))
        self._img(slide, buf, Inches(0.5), Inches(1.0), Inches(12.3), Inches(5.5))

    def _slide_appendix_lifting(self, prs, blank, analysis, rc):
        slide = self._add_slide(prs, blank)
        self._header(slide, '[부록] 리프팅 기기별 세부 분석', 15, rc)

        lift = analysis.get('lifting_detail')
        if lift is None or lift.empty:
            self._text(slide, '리프팅 데이터 없음', Inches(5), Inches(3.5), Inches(3), Inches(1),
                       size=14, color=CG, align=PP_ALIGN.CENTER)
            return

        labels = lift['order_name'].tolist()[:12]
        values = lift['revenue'].tolist()[:12]
        buf = cu.hbar_chart(labels, values, title='리프팅 시술별 매출 TOP 12', figsize=(8, 5))
        self._img(slide, buf, Inches(0.3), Inches(0.9), Inches(12.7), Inches(5.5))

    def _slide_appendix_vip(self, prs, blank, analysis, rc):
        slide = self._add_slide(prs, blank)
        self._header(slide, '[부록] VIP 환자 분석 (상위 10%)', 16, rc)

        pt = analysis.get('patient_spending')
        if pt is None or pt.empty:
            self._text(slide, '데이터 없음', Inches(5), Inches(3.5), Inches(3), Inches(1),
                       size=14, color=CG, align=PP_ALIGN.CENTER)
            return

        threshold = pt['total'].quantile(0.9)
        vip = pt[pt['total'] >= threshold].sort_values('total', ascending=False).head(20)

        self._rect(slide, Inches(0.4), Inches(1.0), Inches(12.5), Inches(0.4), CL)
        self._text(slide, f"상위 10% 기준: {_fmt(threshold)} 이상 | VIP 환자 수: {len(vip):,}명",
                   Inches(0.6), Inches(1.05), Inches(12), Inches(0.3),
                   size=10, color=CP)

        headers = ['차트번호', '총 매출', '방문횟수']
        rows = [[str(r.name), _fmt(r['total']), f"{int(r['visits'])}회"]
                for _, r in vip.iterrows()]
        self._table_text(slide, headers, rows[:18],
                         Inches(0.4), Inches(1.5), Inches(12.5), Inches(5.0))

    def _slide_appendix_misc(self, prs, blank, analysis, rc):
        slide = self._add_slide(prs, blank)
        self._header(slide, '[부록] 기타(미분류) 항목 상세', 17, rc)

        misc = analysis.get('misc_detail')
        if misc is None or misc.empty:
            self._text(slide, '미분류 데이터 없음', Inches(5), Inches(3.5), Inches(3), Inches(1),
                       size=14, color=CG, align=PP_ALIGN.CENTER)
            return

        headers = ['항목명', '매출', '건수']
        rows = [[r['order_name'], _fmt(r['revenue']), f"{int(r['count']):,}건"]
                for _, r in misc.iterrows()]
        self._table_text(slide, headers, rows[:18],
                         Inches(0.4), Inches(1.0), Inches(12.5), Inches(5.5))
